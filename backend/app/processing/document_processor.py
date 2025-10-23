import hashlib
from typing import List, Dict, Any, Optional, Tuple
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import logging
from anthropic import Anthropic
from dotenv import load_dotenv
import os
import re

load_dotenv()

logger = logging.getLogger(__name__)

class DocumentProcessor:
    def __init__(self, chunk_size: int = None, chunk_overlap: int = None, preserve_lists: bool = True):
        api_key = os.getenv("ANTHROPIC_API_KEY")
        if api_key:
            self.anthropic = Anthropic(api_key=api_key)
        else:
            self.anthropic = None
            print("WARNING: ANTHROPIC_API_KEY not set, AI features will be disabled")

        # Use provided chunk size or default to 3000 characters (roughly 750 tokens)
        # Converting characters to tokens: ~4 chars per token
        self.chunk_size_chars = chunk_size if chunk_size is not None else 3000
        self.chunk_overlap_chars = chunk_overlap if chunk_overlap is not None else 200

        # Convert to tokens for backward compatibility (rough estimate: 4 chars per token)
        self.chunk_size = self.chunk_size_chars // 4
        self.chunk_overlap = self.chunk_overlap_chars // 4

        # List preservation setting
        self.preserve_lists = preserve_lists
        
    async def process_document(self, file_path: str, file_type: str) -> Dict[str, Any]:
        """Process document and extract hierarchical structure"""
        
        # Extract text based on file type - NOW WITH PAGE TRACKING
        if file_type == 'pdf':
            text, page_map = await self.extract_pdf_text_with_pages(file_path)
        elif file_type in ['docx', 'doc']:
            text = await self.extract_docx_text(file_path)
            page_map = None  # Word docs don't have fixed pages
        elif file_type in ['pptx', 'ppt']:
            text, page_map = await self.extract_pptx_text_with_slides(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
        
        # Generate document ID
        doc_id = hashlib.md5(text.encode()).hexdigest()
        
        # Generate document summary
        summary = await self.generate_summary(text)
        
        # Extract key concepts
        concepts = await self.extract_concepts(text)
        
        # Identify sections
        sections = await self.identify_sections(text)
        
        # Create chunks with overlap AND PAGE NUMBERS
        chunks = await self.create_smart_chunks(text, sections, page_map)
        
        return {
            'doc_id': doc_id,
            'text': text,
            'summary': summary,
            'concepts': concepts,
            'sections': sections,
            'chunks': chunks,
            'file_type': file_type,
            'file_path': file_path
        }
    
    async def extract_pdf_text_with_pages(self, file_path: str) -> Tuple[str, List[Dict]]:
        """Extract text from PDF with page number tracking"""
        try:
            reader = PdfReader(file_path)
            full_text = ""
            page_map = []  # Track where each page starts/ends in the text
            
            current_position = 0
            for page_num, page in enumerate(reader.pages, start=1):
                page_text = page.extract_text()
                if page_text:
                    start_pos = current_position
                    end_pos = current_position + len(page_text)
                    
                    page_map.append({
                        'page_number': page_num,
                        'start_char': start_pos,
                        'end_char': end_pos,
                        'text': page_text
                    })
                    
                    full_text += page_text + "\n"
                    current_position = end_pos + 1  # +1 for the newline
            
            return full_text, page_map
        except Exception as e:
            logger.error(f"Error extracting PDF text: {e}")
            raise
    
    async def extract_pptx_text_with_slides(self, file_path: str) -> Tuple[str, List[Dict]]:
        """Extract text from PowerPoint with slide number tracking"""
        try:
            prs = Presentation(file_path)
            full_text = ""
            page_map = []  # Treat slides as pages
            
            current_position = 0
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                
                if slide_text:
                    start_pos = current_position
                    end_pos = current_position + len(slide_text)
                    
                    page_map.append({
                        'page_number': slide_num,
                        'start_char': start_pos,
                        'end_char': end_pos,
                        'text': slide_text
                    })
                    
                    full_text += slide_text
                    current_position = end_pos
            
            return full_text, page_map
        except Exception as e:
            logger.error(f"Error extracting PPTX text: {e}")
            raise
    
    async def extract_docx_text(self, file_path: str) -> str:
        """Extract text from Word document"""
        try:
            doc = DocxDocument(file_path)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            return text
        except Exception as e:
            logger.error(f"Error extracting DOCX text: {e}")
            raise
    
    def find_page_for_position(self, char_position: int, page_map: List[Dict]) -> Optional[int]:
        """Find which page a character position belongs to"""
        if not page_map:
            return None

        for page_info in page_map:
            if page_info['start_char'] <= char_position <= page_info['end_char']:
                return page_info['page_number']

        # If position is beyond last page, return last page number
        if page_map and char_position > page_map[-1]['end_char']:
            return page_map[-1]['page_number']

        return None

    def _is_list_item(self, line: str) -> bool:
        """Check if a line is a list item (numbered or bulleted)"""
        stripped = line.strip()
        if not stripped:
            return False

        # Numbered lists: "1.", "2)", "1.1", "a.", "i.", etc.
        numbered_patterns = [
            r'^\d+\.',  # 1. 2. 3.
            r'^\d+\)',  # 1) 2) 3)
            r'^\d+\.\d+',  # 1.1, 1.2, 2.1
            r'^[a-z]\.',  # a. b. c.
            r'^[A-Z]\.',  # A. B. C.
            r'^[ivxlcdm]+\.',  # i. ii. iii. (Roman numerals)
        ]

        # Bulleted lists: "•", "-", "*", "◦", "▪", "▫"
        bulleted_chars = ['•', '-', '*', '◦', '▪', '▫', '–', '—', '●', '○']

        # Check numbered patterns
        for pattern in numbered_patterns:
            if re.match(pattern, stripped):
                return True

        # Check bulleted patterns (must be at start of line after whitespace)
        if stripped[0] in bulleted_chars:
            # Make sure it's not just a dash in the middle of text
            # It should be followed by whitespace
            if len(stripped) > 1 and stripped[1] in [' ', '\t']:
                return True

        return False

    def _detect_list_boundaries(self, text: str) -> List[Tuple[int, int]]:
        """Detect start and end positions of lists in text
        Returns list of (start_char, end_char) tuples"""
        lines = text.split('\n')
        list_boundaries = []

        in_list = False
        list_start_line = 0
        list_start_char = 0

        current_char = 0
        for i, line in enumerate(lines):
            line_start_char = current_char

            if self._is_list_item(line):
                if not in_list:
                    # Start of a new list
                    in_list = True
                    list_start_line = i
                    list_start_char = line_start_char
            else:
                if in_list:
                    # Check if this is just a continuation line (indented)
                    # or the end of the list
                    if line.strip() and not line.startswith(' ') and not line.startswith('\t'):
                        # End of list (non-empty, non-indented line)
                        list_end_char = line_start_char
                        list_boundaries.append((list_start_char, list_end_char))
                        in_list = False
                    # If empty line or indented, continue the list

            current_char += len(line) + 1  # +1 for newline

        # If we ended while still in a list
        if in_list:
            list_boundaries.append((list_start_char, current_char))

        return list_boundaries

    def _is_inside_list(self, char_position: int, list_boundaries: List[Tuple[int, int]]) -> bool:
        """Check if a character position falls inside a list"""
        for start, end in list_boundaries:
            if start <= char_position <= end:
                return True
        return False

    def _find_safe_chunk_boundary(self, text: str, ideal_position: int, list_boundaries: List[Tuple[int, int]],
                                   search_backwards: bool = True) -> int:
        """Find a safe position for chunk boundary that doesn't split lists

        Args:
            text: The full text
            ideal_position: The ideal chunk boundary position
            list_boundaries: List of (start, end) tuples for detected lists
            search_backwards: If True, search backwards for boundary; if False, search forwards

        Returns:
            A safe character position that doesn't split a list
        """
        # If we're not in a list, return the ideal position
        if not self._is_inside_list(ideal_position, list_boundaries):
            return ideal_position

        # We're inside a list, need to find the list boundary
        for start, end in list_boundaries:
            if start <= ideal_position <= end:
                if search_backwards:
                    # Move boundary to before the list starts
                    return start
                else:
                    # Move boundary to after the list ends
                    return end

        return ideal_position
    
    async def create_smart_chunks(self, text: str, sections: List[Dict], page_map: Optional[List[Dict]] = None) -> List[Dict[str, Any]]:
        """Create overlapping chunks with section awareness, PAGE NUMBERS, and LIST PRESERVATION"""
        chunks = []
        chunk_id = 0

        for section in sections:
            section_text = section['text']
            section_start_char = text.find(section_text)  # Find where this section starts in full text

            # Detect list boundaries in this section if preserve_lists is enabled
            list_boundaries = []
            if self.preserve_lists:
                list_boundaries = self._detect_list_boundaries(section_text)
                if list_boundaries:
                    logger.info(f"Detected {len(list_boundaries)} lists in section '{section['title']}'")

            # Use character-based chunking instead of word-based to better handle list boundaries
            target_chunk_chars = self.chunk_size_chars
            overlap_chars = self.chunk_overlap_chars

            start_char = 0
            while start_char < len(section_text):
                # Calculate ideal end position
                end_char = min(start_char + target_chunk_chars, len(section_text))

                # If preserve_lists is enabled, adjust chunk boundary to avoid splitting lists
                if self.preserve_lists and list_boundaries and end_char < len(section_text):
                    # Check if the end position falls inside a list
                    if self._is_inside_list(end_char, list_boundaries):
                        # Adjust boundary to not split the list
                        safe_end = self._find_safe_chunk_boundary(
                            section_text,
                            end_char,
                            list_boundaries,
                            search_backwards=True
                        )

                        # If moving backwards would make chunk too small, try moving forwards
                        if safe_end - start_char < target_chunk_chars * 0.5:
                            safe_end = self._find_safe_chunk_boundary(
                                section_text,
                                end_char,
                                list_boundaries,
                                search_backwards=False
                            )

                        end_char = safe_end

                # Extract chunk text
                chunk_text = section_text[start_char:end_char].strip()

                # Skip empty chunks
                if not chunk_text:
                    break

                # Find the character position of this chunk in the full text
                chunk_char_position = section_start_char + start_char if section_start_char >= 0 else 0

                # Find which page this chunk belongs to
                page_number = self.find_page_for_position(chunk_char_position, page_map) if page_map else None

                chunks.append({
                    'chunk_id': f"chunk_{chunk_id}",
                    'section_title': section['title'],
                    'text': chunk_text,
                    'start_char': start_char,
                    'end_char': end_char,
                    'page_number': page_number
                })

                chunk_id += 1

                # Move to next chunk with overlap
                start_char = end_char - overlap_chars

                # Prevent infinite loop - if we're not making progress, break
                if start_char >= len(section_text) or end_char >= len(section_text):
                    break

        return chunks
    
    async def generate_summary(self, text: str) -> str:
        """Generate document summary using Claude"""
        try:
            # Limit text length for summary
            max_chars = 10000
            truncated_text = text[:max_chars] if len(text) > max_chars else text
            
            response = self.anthropic.messages.create(
                model="claude-3-haiku-20240307",
                max_tokens=500,
                messages=[{
                    "role": "user",
                    "content": f"""Provide a comprehensive summary of this document in 2-3 paragraphs. 
                    Focus on the main topics, key concepts, and practical applications.
                    
                    Document text:
                    {truncated_text}"""
                }]
            )
            return response.content[0].text
        except Exception as e:
            logger.error(f"Error generating summary: {e}")
            return "Summary generation failed"
    
    async def extract_concepts(self, text: str) -> List[str]:
        """Extract key concepts from document"""
        try:
            # Simple keyword extraction for now
            # Could enhance with Claude or NLP libraries
            concepts = []
            
            # Look for capitalized phrases (likely important terms)
            pattern = r'\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*\b'
            matches = re.findall(pattern, text)
            
            # Get unique concepts
            seen = set()
            for match in matches:
                if len(match) > 3 and match not in seen:
                    concepts.append(match)
                    seen.add(match)
                    if len(concepts) >= 20:
                        break
            
            return concepts
        except Exception as e:
            logger.error(f"Error extracting concepts: {e}")
            return []
    
    async def identify_sections(self, text: str) -> List[Dict[str, Any]]:
        """Identify document sections"""
        sections = []
        
        # Look for headers (simple approach)
        lines = text.split('\n')
        current_section = ""
        section_start = 0
        
        for i, line in enumerate(lines):
            # Simple heuristic: short lines in caps or with numbers
            if len(line) < 100 and (line.isupper() or re.match(r'^\d+\.', line)):
                if current_section:
                    sections.append({
                        'title': current_section,
                        'start': section_start,
                        'end': i,
                        'text': '\n'.join(lines[section_start:i])
                    })
                current_section = line.strip()
                section_start = i
        
        # Add last section
        if current_section:
            sections.append({
                'title': current_section,
                'start': section_start,
                'end': len(lines),
                'text': '\n'.join(lines[section_start:])
            })
        
        # If no sections found, treat whole document as one section
        if not sections:
            sections.append({
                'title': 'Full Document',
                'start': 0,
                'end': len(lines),
                'text': text
            })
        
        return sections

# Create processor instance
processor = DocumentProcessor()